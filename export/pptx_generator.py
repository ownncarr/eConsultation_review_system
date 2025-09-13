# export/pptx_generator.py
"""
Generate a simple PPTX summary slide deck.

Primary dependency: python-pptx
Fallback: write a plain text .txt if python-pptx is missing.
"""

from pathlib import Path
from typing import Any, Dict, List, Optional
import yaml
import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    _PPTX_AVAILABLE = True
except Exception:
    _PPTX_AVAILABLE = False


def _load_reports_dir(default: str = "reports") -> Path:
    try:
        with open("configs/settings.yaml", "r") as f:
            cfg = yaml.safe_load(f) or {}
            path = cfg.get("paths", {}).get("reports_dir", default)
    except Exception:
        path = default
    p = Path(path)
    p.mkdir(parents=True, exist_ok=True)
    return p


def generate_pptx_report(
    title: str,
    items: List[Dict[str, Any]],
    filename: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Generate a PPTX file summarizing items.

    Args:
        title: Title for the deck.
        items: List of dicts, each dict describing a single analysis entry.
        filename: Optional output filename (without path).

    Returns:
        dict with status and path.
    """
    reports_dir = _load_reports_dir()
    if filename is None:
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
    out_path = reports_dir / filename

    if not _PPTX_AVAILABLE:
        text_path = out_path.with_suffix(".txt")
        try:
            with open(text_path, "w", encoding="utf-8") as f:
                f.write(f"{title}\nGenerated: {datetime.datetime.now().isoformat()}\n\n")
                for item in items:
                    for k, v in item.items():
                        f.write(f"{k}: {v}\n")
                    f.write("\n" + ("-" * 60) + "\n\n")
            return {"ok": True, "path": str(text_path), "error": None, "fallback": True}
        except Exception as e:
            return {"ok": False, "path": str(text_path), "error": str(e), "fallback": True}

    try:
        prs = Presentation()
        # title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_box = slide.shapes.title
        subtitle_box = slide.placeholders[1]
        title_box.text = title
        subtitle_box.text = f"Generated: {datetime.datetime.now().isoformat()}"

        # create a slide per item (or grouped every N items)
        for item in items:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # title & content layout
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]

            title_text = f"Row: {item.get('row', 'n/a')} - {item.get('sentiment_label', '')}"
            title_shape.text = title_text

            # build the content text
            parts = []
            if "summary" in item:
                parts.append(f"Summary: {item.get('summary')}")
            if "original" in item:
                parts.append(f"Original: {str(item.get('original'))[:300]}")
            if "keywords" in item and item.get("keywords"):
                kws = item.get("keywords")
                if isinstance(kws[0], (list, tuple)):
                    kw_str = ", ".join([str(k[0]) for k in kws])
                else:
                    kw_str = ", ".join([str(k) for k in kws])
                parts.append(f"Keywords: {kw_str}")

            tf = body_shape.text_frame
            tf.clear()
            # Add bullet lines
            for p in parts:
                p = p.replace("\n", " ")
                tf.add_paragraph().text = p

        prs.save(str(out_path))
        return {"ok": True, "path": str(out_path), "error": None, "fallback": False}
    except Exception as e:
        return {"ok": False, "path": str(out_path), "error": str(e), "fallback": False}
